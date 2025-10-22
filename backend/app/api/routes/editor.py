"""Editor API routes for managing translation edits."""
from fastapi import APIRouter, HTTPException, Depends
from pydantic import BaseModel, Field
from typing import List
import logging

from app.services.translation_processor import TranslationProcessor
from app.api.dependencies import get_translation_processor
from app.models.translation import ImproveTranslationRequest, ImproveTranslationResponse

logger = logging.getLogger(__name__)
router = APIRouter()

# Cache for slide images to avoid regenerating
_slide_image_cache = {}


class SlideContent(BaseModel):
    """Model for a slide's content."""
    slide_number: int = Field(..., description="Slide number (1-indexed)")
    text_frames: List[dict] = Field(..., description="Text frames in the slide")


class DocumentContentResponse(BaseModel):
    """Response model for document content extraction."""
    success: bool = Field(..., description="Whether extraction succeeded")
    filename: str = Field(..., description="Document filename")
    total_slides: int = Field(..., description="Total number of slides")
    slides: List[SlideContent] = Field(..., description="Slide contents")
    error: str | None = Field(None, description="Error message if failed")


class TranslationEdit(BaseModel):
    """Model for editing a single translation."""
    original_text: str = Field(..., description="Original source text")
    current_translation: str = Field(..., description="Current translation")
    edited_translation: str = Field(..., description="User's edited translation")
    target_language: str = Field(..., description="Target language code")


class BulkEditRequest(BaseModel):
    """Model for bulk editing translations."""
    edits: List[TranslationEdit] = Field(..., description="List of translation edits")


class BulkEditResponse(BaseModel):
    """Response for bulk edit operation."""
    success: bool = Field(..., description="Whether operation succeeded")
    edited_count: int = Field(..., description="Number of translations edited")
    edits: List[TranslationEdit] = Field(..., description="Edited translations")


@router.post("/save-edits", response_model=BulkEditResponse)
async def save_translation_edits(request: BulkEditRequest):
    """
    Save user's manual edits to translations.
    
    This endpoint stores the user's edited translations.
    In a production system, these would be saved to a database.
    
    Args:
        request: Bulk edit request with list of edits
        
    Returns:
        Confirmation of saved edits
    """
    try:
        if not request.edits:
            raise HTTPException(status_code=400, detail="No edits provided")
        
        # In a real application, you would save these to a database
        # For now, we just return them as confirmation
        logger.info(f"Saving {len(request.edits)} translation edits")
        
        return BulkEditResponse(
            success=True,
            edited_count=len(request.edits),
            edits=request.edits
        )
    
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error saving edits: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to save edits: {str(e)}")


@router.get("/slide-preview/{filename}/{slide_number}")
async def get_slide_preview(filename: str, slide_number: int):
    """
    Generate a placeholder preview image (Aspose.Slides disabled due to ICU issues).
    
    Args:
        filename: Document filename
        slide_number: Slide number (0-indexed)
        
    Returns:
        Placeholder image
    """
    from app.config import settings
    from pathlib import Path
    import io
    from fastapi.responses import Response
    from PIL import Image, ImageDraw
    
    try:
        file_path = settings.OUTPUT_FOLDER / filename
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Check cache first
        file_mtime = file_path.stat().st_mtime
        cache_key = f"{filename}_{slide_number}_{file_mtime}"
        
        if cache_key in _slide_image_cache:
            logger.debug(f"Serving cached placeholder for {cache_key}")
            return Response(
                content=_slide_image_cache[cache_key],
                media_type="image/png",
                headers={
                    "Cache-Control": "public, max-age=86400",
                    "Content-Disposition": f"inline; filename=slide_{slide_number}.png"
                }
            )
        
        # Create placeholder image
        img = Image.new('RGB', (960, 720), color='#f8f9fa')
        draw = ImageDraw.Draw(img)
        
        # Draw slide number and message
        text_lines = [
            f"Slide {slide_number + 1}",
            "",
            "Preview not available",
            "",
            "(Download file to see presentation)"
        ]
        
        y_position = 280
        for line in text_lines:
            bbox = draw.textbbox((0, 0), line)
            text_width = bbox[2] - bbox[0]
            x_position = (960 - text_width) // 2
            draw.text((x_position, y_position), line, fill='#666666')
            y_position += 40
        
        # Draw border
        draw.rectangle([(0, 0), (959, 719)], outline='#dee2e6', width=3)
        
        # Save to bytes
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG')
        img_bytes = img_buffer.getvalue()
        
        # Cache the placeholder
        _slide_image_cache[cache_key] = img_bytes
        if len(_slide_image_cache) > 20:
            oldest_key = next(iter(_slide_image_cache))
            del _slide_image_cache[oldest_key]
        
        return Response(
            content=img_bytes,
            media_type="image/png",
            headers={
                "Cache-Control": "public, max-age=86400",
                "Content-Disposition": f"inline; filename=slide_{slide_number}.png"
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating placeholder: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to generate preview: {str(e)}")


@router.get("/document-content/{filename}", response_model=DocumentContentResponse)
async def get_document_content(filename: str):
    """
    Extract all text content from a translated document for editing.
    
    Args:
        filename: Name of the translated document file
        
    Returns:
        Document content with all slides and text frames
    """
    from app.config import settings
    from pptx import Presentation
    from pathlib import Path
    
    try:
        file_path = settings.OUTPUT_FOLDER / filename
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Document not found")
        
        prs = Presentation(file_path)
        slides_content = []
        
        for slide_idx, slide in enumerate(prs.slides):
            text_frames = []
            
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    text_frames.append({
                        'id': f'slide_{slide_idx}_shape_{shape_idx}',
                        'text': shape.text_frame.text,
                        'shape_index': shape_idx
                    })
            
            if text_frames:  # Only include slides with text
                slides_content.append(SlideContent(
                    slide_number=slide_idx + 1,
                    text_frames=text_frames
                ))
        
        return DocumentContentResponse(
            success=True,
            filename=filename,
            total_slides=len(prs.slides),
            slides=slides_content
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error extracting document content: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract content: {str(e)}")


class UpdateContentRequest(BaseModel):
    """Request model for updating document content."""
    filename: str = Field(..., description="Document filename to update")
    edits: List[dict] = Field(..., description="List of text edits with id and new text")


@router.post("/preview-with-edits")
async def preview_slide_with_edits(request: UpdateContentRequest, slide_number: int):
    """
    Generate a preview of a slide with temporary edits applied (without saving).
    
    Args:
        request: Edits to apply temporarily
        slide_number: Slide number to preview
        
    Returns:
        Image of the slide with edits applied
    """
    from app.config import settings
    from pathlib import Path
    import io
    from fastapi.responses import StreamingResponse
    import aspose.slides as slides
    from pptx import Presentation
    import tempfile
    
    try:
        file_path = settings.OUTPUT_FOLDER / request.filename
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Create a temporary copy
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            tmp_path = Path(tmp.name)
        
        # Load and modify
        prs = Presentation(file_path)
        edits_map = {edit['id']: edit['text'] for edit in request.edits}
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    frame_id = f'slide_{slide_idx}_shape_{shape_idx}'
                    if frame_id in edits_map:
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                        p.text = edits_map[frame_id]
        
        prs.save(tmp_path)
        
        # Generate image from temp file
        with slides.Presentation(str(tmp_path)) as prs_slides:
            if slide_number < 0 or slide_number >= len(prs_slides.slides):
                raise HTTPException(status_code=400, detail="Invalid slide number")
            
            slide = prs_slides.slides[slide_number]
            scale = 2.0
            with slide.get_image(scale, scale) as img:
                img_buffer = io.BytesIO()
                img.save(img_buffer, slides.ImageFormat.PNG)
                img_buffer.seek(0)
                
                # Clean up temp file
                tmp_path.unlink()
                
                return StreamingResponse(
                    img_buffer,
                    media_type="image/png",
                    headers={"Cache-Control": "no-cache"}
                )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error generating preview: {e}")
        if tmp_path.exists():
            tmp_path.unlink()
        raise HTTPException(status_code=500, detail=f"Failed to generate preview: {str(e)}")


@router.post("/update-content")
async def update_document_content(request: UpdateContentRequest):
    """
    Apply user's edits back to the document.
    
    Args:
        request: Update request with filename and list of edits
        
    Returns:
        Success response with updated filename
    """
    from app.config import settings
    from pptx import Presentation
    from pathlib import Path
    
    try:
        file_path = settings.OUTPUT_FOLDER / request.filename
        if not file_path.exists():
            raise HTTPException(status_code=404, detail="Document not found")
        
        prs = Presentation(file_path)
        
        # Create a mapping of edits by id
        edits_map = {edit['id']: edit['text'] for edit in request.edits}
        
        # Apply edits
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    frame_id = f'slide_{slide_idx}_shape_{shape_idx}'
                    if frame_id in edits_map:
                        new_text = edits_map[frame_id]
                        # Clear all existing paragraphs and create new ones
                        text_frame = shape.text_frame
                        text_frame.clear()
                        # Add the new text
                        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                        p.text = new_text
        
        # Save the updated presentation
        prs.save(file_path)
        
        # Clear cache for this file
        global _slide_image_cache
        file_mtime = file_path.stat().st_mtime
        keys_to_delete = [k for k in _slide_image_cache.keys() if k.startswith(f"{request.filename}_")]
        for key in keys_to_delete:
            del _slide_image_cache[key]
        
        logger.info(f"Document updated with {len(request.edits)} edits: {request.filename}")
        
        return {
            "success": True,
            "filename": request.filename,
            "edits_applied": len(request.edits)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error updating document: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to update document: {str(e)}")


@router.post("/suggest-improvement", response_model=ImproveTranslationResponse)
async def suggest_translation_improvement(
    request: ImproveTranslationRequest,
    processor: TranslationProcessor = Depends(get_translation_processor)
):
    """
    Get AI suggestions for improving a translation.
    
    Uses LLM to suggest improvements to the current translation.
    
    Args:
        request: Improvement request with original and current translation
        processor: Translation processor instance
        
    Returns:
        Suggested improved translation
    """
    try:
        result = processor.improve_translation(
            original_text=request.original_text,
            current_translation=request.current_translation,
            target_language=request.target_language,
            feedback=request.feedback
        )
        
        return ImproveTranslationResponse(**result)
    
    except Exception as e:
        logger.error(f"Error suggesting improvement: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to suggest improvement: {str(e)}")