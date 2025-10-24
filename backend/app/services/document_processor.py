"""
Document Processor Service for PPTX Translation.

This service handles:
1. PPTX text extraction from slides, shapes, and tables.
2. Formatting preservation during translation.
3. Integration with Translation Processor for Azure and OpenRouter translation.
4. OCR-based translation of text embedded in images.

Algorithm:
- Extracts text from PPTX slides and shapes
- Extracts and translates text from images using OCR
- Translates text while preserving formatting
- Creates new translated PPTX file
"""

import logging
import hashlib
from pathlib import Path
from typing import Any, Dict, Optional
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import io

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Processes PPTX documents for translation with formatting preservation."""

    def __init__(self, translation_processor, image_translator=None):
        """
        Initialize the DocumentProcessor.

        Args:
            translation_processor: An instance of the translation processor to handle translation logic.
            image_translator: An instance of the image translator for OCR-based image translation (optional).
        """
        self.translation_processor = translation_processor
        self.image_translator = image_translator
        self.original_texts = {}  # Store original texts for before/after comparison
        logger.info("DocumentProcessor initialized")
        if image_translator:
            logger.info("Image translation enabled")

    def process_pptx(
        self,
        input_path: Path,
        output_path: Path,
        target_language: str,
        source_language: Optional[str] = None,
        use_llm: bool = False,
        llm_model: Optional[str] = None,
        preserve_formatting: bool = True
    ) -> Dict[str, Any]:
        """
        Process PPTX file and create translated version.

        Args:
            input_path: Path to input PPTX file
            output_path: Path to save translated PPTX
            target_language: Target language code
            source_language: Source language code (optional)
            use_llm: Whether to use LLM enhancement
            llm_model: LLM model to use (optional)
            preserve_formatting: Whether to preserve original formatting

        Returns:
            Dictionary with processing statistics
        """
        logger.info(f"Processing PPTX: {input_path.name}")
        
        stats = {
            'filename': input_path.name,
            'slides_processed': 0,
            'text_frames_translated': 0,
            'tables_translated': 0,
            'images_translated': 0,
            'source_language': source_language,
            'target_language': target_language,
            'method': 'llm' if use_llm else 'azure'
        }

        try:
            # Reset original texts storage
            self.original_texts = {}
            
            # Load presentation
            prs = Presentation(input_path)
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                logger.info(f"Processing slide {slide_idx + 1}/{len(prs.slides)}")
                
                # Create a fixed list of shapes to avoid modifying collection during iteration
                # (image replacement adds new shapes which would cause infinite loop)
                shapes_to_process = list(slide.shapes)
                
                # Track processed images by their binary content to avoid duplicates
                processed_image_hashes = set()
                
                for shape_idx, shape in enumerate(shapes_to_process):
                    # Process GROUP shapes recursively (they contain nested shapes)
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        self._process_group_shape(
                            shape,
                            slide,
                            target_language,
                            source_language,
                            use_llm,
                            llm_model,
                            preserve_formatting,
                            slide_idx,
                            shape_idx,
                            stats,
                            processed_image_hashes
                        )
                        continue
                    
                    # Process images with text (OCR translation) FIRST
                    # This must come before text frame processing to avoid conflicts
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.image_translator:
                        # Get image hash to check for duplicates
                        try:
                            image_hash = hashlib.md5(shape.image.blob).hexdigest()
                            
                            # Skip if we've already processed this exact image
                            if image_hash in processed_image_hashes:
                                continue
                            
                            processed_image_hashes.add(image_hash)
                        except:
                            pass  # If hashing fails, continue anyway
                        
                        if self._process_image(
                            shape,
                            slide,
                            target_language,
                            source_language,
                            use_llm,
                            llm_model
                        ):
                            stats['images_translated'] += 1
                        
                        # After processing image, skip to next shape
                        # (don't process as text frame even if it has one)
                        continue
                    
                    # Process text frames (text boxes, titles, etc.)
                    if shape.has_text_frame:
                        self._process_text_frame(
                            shape.text_frame,
                            target_language,
                            source_language,
                            use_llm,
                            llm_model,
                            preserve_formatting,
                            slide_idx,
                            shape_idx
                        )
                        stats['text_frames_translated'] += 1
                    
                    # Process tables
                    if shape.has_table:
                        self._process_table(
                            shape.table,
                            target_language,
                            source_language,
                            use_llm,
                            llm_model
                        )
                        stats['tables_translated'] += 1
                
                stats['slides_processed'] += 1
            
            # Save translated presentation
            prs.save(output_path)
            logger.info(f"Translated PPTX saved to: {output_path}")
            
            # Save original texts mapping to a JSON file
            original_texts_path = output_path.with_suffix('.original.json')
            with open(original_texts_path, 'w', encoding='utf-8') as f:
                json.dump(self.original_texts, f, ensure_ascii=False, indent=2)
            logger.info(f"Original texts saved to: {original_texts_path}")
            
            return {
                'success': True,
                **stats
            }
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            raise

    def _process_group_shape(
        self,
        group_shape,
        slide,
        target_language: str,
        source_language: Optional[str],
        use_llm: bool,
        llm_model: Optional[str],
        preserve_formatting: bool,
        slide_idx: int,
        parent_shape_idx: int,
        stats: Dict,
        processed_image_hashes: set
    ):
        """
        Recursively process shapes within a group.
        Groups can contain text boxes, images, and even nested groups.
        """
        try:
            for nested_idx, nested_shape in enumerate(group_shape.shapes):
                # Recursively process nested groups
                if nested_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._process_group_shape(
                        nested_shape,
                        slide,
                        target_language,
                        source_language,
                        use_llm,
                        llm_model,
                        preserve_formatting,
                        slide_idx,
                        parent_shape_idx,
                        stats,
                        processed_image_hashes
                    )
                    continue
                
                # Process images in group
                if nested_shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.image_translator:
                    try:
                        image_hash = hashlib.md5(nested_shape.image.blob).hexdigest()
                        if image_hash in processed_image_hashes:
                            continue
                        processed_image_hashes.add(image_hash)
                    except:
                        pass
                    
                    if self._process_image(
                        nested_shape,
                        slide,
                        target_language,
                        source_language,
                        use_llm,
                        llm_model
                    ):
                        stats['images_translated'] += 1
                    continue
                
                # Process text frames in group
                if nested_shape.has_text_frame:
                    self._process_text_frame(
                        nested_shape.text_frame,
                        target_language,
                        source_language,
                        use_llm,
                        llm_model,
                        preserve_formatting,
                        slide_idx,
                        f"{parent_shape_idx}_group_{nested_idx}"
                    )
                    stats['text_frames_translated'] += 1
                
                # Process tables in group
                if nested_shape.has_table:
                    self._process_table(
                        nested_shape.table,
                        target_language,
                        source_language,
                        use_llm,
                        llm_model
                    )
                    stats['tables_translated'] += 1
        except Exception as e:
            logger.error(f"Error processing group shape: {e}")


        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return {
                'success': False,
                'error': str(e),
                **stats
            }

    def _process_text_frame(
        self,
        text_frame,
        target_language: str,
        source_language: Optional[str],
        use_llm: bool,
        llm_model: Optional[str],
        preserve_formatting: bool,
        slide_idx: int = 0,
        shape_idx: int = 0
    ):
        """Process a text frame and translate its content."""
        try:
            original_text = text_frame.text.strip()
            if not original_text:
                return
            
            logger.info(f"Translating text frame on slide {slide_idx + 1}: '{original_text[:50]}...'")
            
            # Create unique ID for this text frame
            frame_id = f'slide_{slide_idx}_shape_{shape_idx}'
            
            # Store original text
            self.original_texts[frame_id] = original_text
            
            # Translate text
            result = self.translation_processor.translate_text(
                text=original_text,
                target_language=target_language,
                source_language=source_language,
                force_llm=use_llm,
                llm_model=llm_model
            )
            
            if result.get('success') and result.get('translation'):
                translated_text = result['translation']
                logger.info(f"Translated text frame: '{original_text[:30]}' -> '{translated_text[:30]}'")
                
                if preserve_formatting:
                    self._replace_text_preserve_format(text_frame, translated_text)
                else:
                    text_frame.text = translated_text
            else:
                logger.warning(f"Translation failed for text frame: {result.get('error', 'Unknown error')}")
                    
        except Exception as e:
            logger.error(f"Error processing text frame: {e}")

    def _process_table(
        self,
        table,
        target_language: str,
        source_language: Optional[str],
        use_llm: bool,
        llm_model: Optional[str]
    ):
        """Process a table and translate its cells."""
        try:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        result = self.translation_processor.translate_text(
                            text=cell.text.strip(),
                            target_language=target_language,
                            source_language=source_language,
                            force_llm=use_llm,
                            llm_model=llm_model
                        )
                        
                        if result.get('success') and result.get('translation'):
                            cell.text = result['translation']
                            
        except Exception as e:
            logger.error(f"Error processing table: {e}")

    def _process_image(
        self,
        shape,
        slide,
        target_language: str,
        source_language: Optional[str],
        use_llm: bool,
        llm_model: Optional[str]
    ) -> bool:
        """
        Process an image shape and translate embedded text using OCR.
        
        Args:
            shape: Picture shape from slide
            slide: Parent slide object
            target_language: Target language code
            source_language: Source language code (optional)
            use_llm: Whether to use LLM enhancement
            llm_model: LLM model to use
            
        Returns:
            True if image was successfully translated, False otherwise
        """
        try:
            # Get image from shape
            image = shape.image
            image_bytes = image.blob
            content_type = image.content_type
            
            # Skip very small images (likely decorative icons)
            if len(image_bytes) < 5000:  # Less than 5KB
                logger.debug(f"Skipping small image ({len(image_bytes)} bytes), likely decorative")
                return False
            
            logger.info(f"Processing image: {content_type}, size: {len(image_bytes)} bytes, dimensions: {shape.width} x {shape.height}")
            
            # Translate the image
            translated_image_bytes = self.image_translator.translate_image(
                image_bytes=image_bytes,
                content_type=content_type,
                translation_processor=self.translation_processor,
                target_language=target_language,
                source_language=source_language,
                use_llm=use_llm,
                llm_model=llm_model
            )
            
            if translated_image_bytes:
                # Replace the image in the slide
                # Get shape properties
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Remove old shape
                sp = shape.element
                sp.getparent().remove(sp)
                
                # Add new image with translated text
                pic = slide.shapes.add_picture(
                    io.BytesIO(translated_image_bytes),
                    left, top, width, height
                )
                
                logger.info("Image successfully translated and replaced")
                return True
            else:
                logger.info("No text found in image or translation skipped")
                return False
                
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            return False

    def _replace_text_preserve_format(self, text_frame, new_text: str):
        """
        Replace text in a text frame while preserving formatting.

        Args:
            text_frame: The text frame object from the document.
            new_text: The new text to insert.
        """
        try:
            # Preserve first paragraph's formatting
            if text_frame.paragraphs:
                first_para = text_frame.paragraphs[0]
                if first_para.runs:
                    first_run = first_para.runs[0]
                    font_size = first_run.font.size
                    font_name = first_run.font.name
                    font_bold = first_run.font.bold
                    font_italic = first_run.font.italic
                    
                    # Clear and replace
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = new_text
                    
                    # Apply preserved formatting
                    if font_size:
                        run.font.size = font_size
                    if font_name:
                        run.font.name = font_name
                    if font_bold is not None:
                        run.font.bold = font_bold
                    if font_italic is not None:
                        run.font.italic = font_italic
                else:
                    text_frame.text = new_text
            else:
                text_frame.text = new_text
                
        except Exception as e:
            logger.error(f"Error preserving format: {e}")
            text_frame.text = new_text